"""Document processing pipeline for various file formats."""
import os
import hashlib
from pathlib import Path
from typing import List, Dict, Optional
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
import markdown
from bs4 import BeautifulSoup
import chardet

from utils.config_loader import ConfigLoader
from utils.memory_monitor import MemoryMonitor


class DocumentProcessor:
    """Process various document formats and extract text."""
    
    def __init__(self, config: ConfigLoader):
        self.config = config
        self.memory_monitor = MemoryMonitor(
            max_memory_mb=config.get('memory.max_memory_usage_mb', 6000)
        )
        self.supported_formats = config.get('document.supported_formats', [])
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from a document based on its format."""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        try:
            if file_ext == '.pdf':
                return self._extract_from_pdf(file_path)
            elif file_ext == '.docx':
                return self._extract_from_docx(file_path)
            elif file_ext == '.pptx':
                return self._extract_from_pptx(file_path)
            elif file_ext == '.txt':
                return self._extract_from_txt(file_path)
            elif file_ext == '.md':
                return self._extract_from_markdown(file_path)
            elif file_ext == '.xlsx':
                return self._extract_from_excel(file_path)
            else:
                raise ValueError(f"Unsupported format: {file_ext}")
        except Exception as e:
            raise Exception(f"Error extracting text from {file_path}: {str(e)}")
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file using PyMuPDF (fitz)."""
        text = []
        try:
            with fitz.open(file_path) as doc:
                for page in doc:
                    text.append(page.get_text())
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"PyMuPDF error: {e}")
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs]
        return '\n'.join(paragraphs)
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file with encoding detection."""
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            encoding = chardet.detect(raw_data)['encoding']
        
        with open(file_path, 'r', encoding=encoding or 'utf-8') as f:
            return f.read()
    
    def _extract_from_markdown(self, file_path: str) -> str:
        """Extract text from Markdown file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, 'html.parser')
        return soup.get_text()
    
    def _extract_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file."""
        workbook = openpyxl.load_workbook(file_path)
        text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) if cell else '' for cell in row)
                if row_text.strip():
                    text.append(row_text)
        return '\n'.join(text)
    
    def chunk_text(self, text: str, chunk_size: Optional[int] = None, 
                   chunk_overlap: Optional[int] = None) -> List[Dict[str, any]]:
        """Split text into chunks with metadata."""
        chunk_size = chunk_size or self.config.get('document.chunk_size', 512)
        chunk_overlap = chunk_overlap or self.config.get('document.chunk_overlap', 50)
        
        chunks = []
        start = 0
        chunk_id = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end]
            
            # Find a good break point (sentence or word boundary)
            if end < len(text):
                # Try to break at sentence end
                last_period = chunk_text.rfind('.')
                last_newline = chunk_text.rfind('\n')
                break_point = max(last_period, last_newline)
                
                if break_point > chunk_size * 0.5:  # Don't break too early
                    chunk_text = chunk_text[:break_point + 1]
                    end = start + break_point + 1
            
            if chunk_text.strip():
                chunks.append({
                    'chunk_id': chunk_id,
                    'text': chunk_text.strip(),
                    'start_pos': start,
                    'end_pos': end,
                    'length': len(chunk_text)
                })
                chunk_id += 1
            
            start = end - chunk_overlap  # Overlap for context
            
            # Memory check
            if not self.memory_monitor.check_memory_available():
                self.memory_monitor.force_gc()
        
        return chunks
    
    def process_document(self, file_path: str) -> Dict[str, any]:
        """Process a document and return chunks with metadata."""
        file_path = Path(file_path)
        
        # Calculate file hash for deduplication
        with open(file_path, 'rb') as f:
            file_hash = hashlib.md5(f.read()).hexdigest()
        
        # Extract text
        text = self.extract_text(str(file_path))
        
        # Chunk text
        chunks = self.chunk_text(text)
        
        return {
            'file_path': str(file_path),
            'file_name': file_path.name,
            'file_hash': file_hash,
            'file_size': file_path.stat().st_size,
            'text_length': len(text),
            'chunks': chunks,
            'chunk_count': len(chunks)
        }



